from datetime import datetime
from io import BytesIO
import requests
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import ChartData
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx import Presentation
from openpyxl.styles.borders import Border, Side, BORDER_THIN
import json
import os
from django.http import HttpResponse, JsonResponse
from django.views import View
from django.views.decorators.csrf import csrf_exempt
from rest_framework.parsers import JSONParser
from openpyxl import Workbook
from openpyxl.chart import Reference, PieChart3D, LineChart

from openpyxl.styles import Font
try:
    from openpyxl.cell import get_column_letter
except ImportError:
    from openpyxl.utils import get_column_letter
    from openpyxl.utils import column_index_from_string
from openpyxl.styles import PatternFill
from openpyxl.styles import Alignment
from openpyxl.drawing.image import Image
import io
import urllib3
from openpyxl import load_workbook


@csrf_exempt
def downloadExcel(request, *args, **kwargs):
    excelData = JSONParser().parse(request)
    # return JsonResponse(excelData, safe=False)
    workBook = Workbook()

    overview = workBook.active
    overview.title = 'Overview'
    workBook = setSingleReportPage(workBook, excelData['Report'])
    overview = setOverviewPage(overview, excelData)

    isExist = os.path.exists('media/estimate/excel/')
    if not isExist:
        os.makedirs('media/estimate/excel/')
    workBook.save('media/estimate/excel/'+excelData['PlanName']+".xlsx")

    return JsonResponse('/media/estimate/excel/'+excelData['PlanName']+".xlsx", safe=False)


def setOverviewPage(overview, excelData):

    http = urllib3.PoolManager()
    r = http.request(
        'GET', 'https://globalmediakit.com/assets/images/logo.png')
    image_file = io.BytesIO(r.data)
    img = Image(image_file)
    img.height = 75
    img.width = 250
    overview.add_image(img, 'A2')
    overview['B5']
    overview.cell(2, 3)
    for sheetRow in excelData['OverviewTable']:
        overview.append(sheetRow)

    min_col = 2
    row_count = len(excelData['AllUnits'])-2
    graph_row_no = 7+row_count+5
    for sheetRow in excelData['AllUnits']:
        # if sheetRow['UnitType']=='Percentage':
        pie = PieChart3D()
        labels = Reference(overview, min_col=1,
                           min_row=7, max_row=7+row_count)
        data = Reference(overview, min_col=min_col,
                         min_row=6, max_row=6+row_count)
        pie.add_data(data, titles_from_data=True)
        pie.set_categories(labels)
        pie.title = sheetRow['CampaignUnit']
        pie.height = 5.6  # default is 7.5
        pie.width = 17  # default is 15
        overview.add_chart(pie, "B"+str(graph_row_no))
        # else:
        #     pie = LineChart()
        #     labels = Reference(overview, min_col=1,
        #             min_row=7, max_row=7+row_count)
        #     data = Reference(overview, min_col=min_col,
        #             min_row=6, max_row=6+row_count)
        #     pie.add_data(data, titles_from_data=True)
        #     pie.set_categories(labels)
        #     pie.title = sheetRow['CampaignUnit']
        #     pie.height = 5.6  # default is 7.5
        #     pie.width = 17  # default is 15
        #     overview.add_chart(pie, "B"+str(graph_row_no))
        min_col = min_col+1
        graph_row_no = graph_row_no+14
    setAutoWidth(overview)
    setSheet(overview)
    return overview


def setSingleReportPage(workBook, excelData):
    workSheetNameCount = 0
    for sheetData in excelData:
        # sheetData = excelData[1]
        workSheetNameCount = workSheetNameCount+1
        workSheetName = "Sheet-"+str(workSheetNameCount)
        # workSheetName = sheetData['SheetName']+str(workSheetNameCount)
        workBook.create_sheet(workSheetName)
        # workBook[workSheetName].merge_cells('A6:B6']
        Cell = workBook[workSheetName]['A5']
        Cell.alignment = Alignment(
            horizontal="left", vertical="center", wrapText=True)
        Cell.font = Font(bold=True, size=12)
        Cell.value = "Media Name"
        # workBook[workSheetName].merge_cells('C6:F6']
        Cell = workBook[workSheetName]['B5']
        Cell.alignment = Alignment(
            horizontal="left", vertical="center", wrapText=True)
        Cell.font = Font(bold=True, size=12)
        Cell.value = sheetData['MediaName']

        # workBook[workSheetName].merge_cells('A7:B7']
        Cell = workBook[workSheetName]['A6']
        Cell.alignment = Alignment(
            horizontal="left", vertical="center", wrapText=True)
        Cell.font = Font(bold=True, size=12)
        Cell.value = "Country/Event"
        # workBook[workSheetName].merge_cells('C7:F7']
        Cell = workBook[workSheetName]['B6']
        Cell.alignment = Alignment(
            horizontal="left", vertical="center", wrapText=True)
        Cell.font = Font(bold=True, size=12)
        Cell.value = sheetData['Country']

        # workBook[workSheetName].merge_cells('A8:B8']
        Cell = workBook[workSheetName]['A8']
        Cell.alignment = Alignment(
            horizontal="left", vertical="center", wrapText=True)
        Cell.font = Font(bold=True, size=12)
        Cell.value = "Language"
        # workBook[workSheetName].merge_cells('C8:F8']
        Cell = workBook[workSheetName]['B8']
        Cell.alignment = Alignment(
            horizontal="left", vertical="center", wrapText=True)
        Cell.font = Font(bold=True, size=12)
        Cell.value = sheetData['Language']

        # workBook[workSheetName].merge_cells('A9:B9']
        Cell = workBook[workSheetName]['A9']
        Cell.alignment = Alignment(
            horizontal="left", vertical="center", wrapText=True)
        Cell.font = Font(bold=True, size=12)
        Cell.value = "Vertical"
        # workBook[workSheetName].merge_cells('C9:F9']
        Cell = workBook[workSheetName]['B9']
        Cell.alignment = Alignment(
            horizontal="left", vertical="center", wrapText=True)
        Cell.font = Font(bold=True, size=12)
        Cell.value = sheetData['VerticalName']

        workBook[workSheetName]['A12']
        workBook[workSheetName].cell(2, 3)
        for sheetRow in sheetData['CampaignReportTable']:
            workBook[workSheetName].append(sheetRow)

        i = 14
        thin_border = Border(
            left=Side(border_style=BORDER_THIN, color='000000'),
            right=Side(border_style=BORDER_THIN, color='000000'),
            top=Side(border_style=BORDER_THIN, color='000000'),
            bottom=Side(border_style=BORDER_THIN, color='000000')
        )
        j = 0
        for sheetRow in sheetData['CampaignReportTable']:
            for cell in workBook[workSheetName][str(i)]:
                cell.alignment = Alignment(
                    horizontal="left", vertical="center", wrapText=True)
                if cell.value or cell.value == 0:
                    cell.border = thin_border
            j = j+1
            if (workBook[workSheetName]['C'+str(i)].value) and workBook[workSheetName]['C'+str(i)].value != 'N/A':
                value = workBook[workSheetName]['C'+str(i)].value
                workBook[workSheetName]['C'+str(i)].value = 'View Proof'
                workBook[workSheetName]['C'+str(i)].hyperlink = value
                workBook[workSheetName]['C' +
                                        str(i)].fill = PatternFill("solid", fgColor="f7941d")
            i = i+1
        min_col = 4
        row_count = len(sheetData['CampaignReportTable'])-1
        graph_row_no = 13+row_count+5
        for sheetRow in sheetData['CampaignReport']:
            # if sheetRow['UnitType'] == 'Percentage':
            pie = PieChart3D()
            # labels = Reference(workBook[workSheetName], min_col=2,
            #        min_row=14, max_row=14+row_count)
            # data = Reference(workBook[workSheetName], min_col=min_col,
            #      min_row=13, max_row=13+row_count)

            piedata = [
                ['Total Served', sheetRow['series'][0]],
                ['Remaining', sheetRow['series'][1]],

            ]
            workBook[workSheetName]["B"+str(graph_row_no+1)]
            workBook[workSheetName].cell(2, 3)
            for d in piedata:
                workBook[workSheetName].append(d)
            labels = Reference(workBook[workSheetName], min_col=1,
                               min_row=graph_row_no+2, max_col=1, max_row=graph_row_no+3)
            data = Reference(workBook[workSheetName], min_col=2,
                             min_row=graph_row_no+1, max_col=2, max_row=graph_row_no+3)
            pie.add_data(data, titles_from_data=True)
            pie.set_categories(labels)
            pie.title = sheetRow['CampaignUnit']
            pie.height = 5.6  # default is 7.5
            pie.width = 17  # default is 15
            workBook[workSheetName].add_chart(pie, "A"+str(graph_row_no))
        # else:
        #     pie = LineChart()
        #     labels = Reference(workBook[workSheetName], min_col=2,
        #                        min_row=14, max_row=14+row_count)
        #     data = Reference(workBook[workSheetName], min_col=min_col,
        #                      min_row=13, max_row=13+row_count)
        #     pie.add_data(data, titles_from_data=True)
        #     pie.set_categories(labels)
        #     pie.title = sheetRow['CampaignUnit']
        #     pie.height = 5.6  # default is 7.5
        #     pie.width = 17  # default is 15
        #     workBook[workSheetName].add_chart(pie, "A"+str(graph_row_no))
            min_col = min_col+1
            graph_row_no = graph_row_no+14
        setAutoWidth(workBook[workSheetName])
        setSheet(workBook[workSheetName])
    return workBook


def setAutoWidth(WorkSheet):
    for column_cells in WorkSheet.columns:
        new_column_length = max(len(str(cell.value))
                                for cell in column_cells)
        new_column_letter = (get_column_letter(column_cells[0].column))
        if new_column_length > 0:
            WorkSheet.column_dimensions[new_column_letter].width = 22
            WorkSheet.column_dimensions[new_column_letter].height = 100
        #     for cell in column_cells:
        #         if cell.value == "Availability":
        #             WorkSheet.column_dimensions[new_column_letter].width = 40
        #         if cell.value == "Media Details":
        #             WorkSheet.column_dimensions[new_column_letter].width = 100
        #         if cell.value == "Media Name":
        #             WorkSheet.column_dimensions[new_column_letter].width = 50
        #         if cell.value == "Countries Covered":
        #             WorkSheet.column_dimensions[new_column_letter].width = 40
        #         if cell.value == "View Image":
        #             cell.font =  Font(bold=True, size=12, color="0056b3")
        for cell in column_cells:
            cell.alignment = Alignment(
                horizontal="left", vertical="center", wrapText=True)


def setSheet(WorkSheet):
    WorkSheet.sheet_view.showGridLines = False
    thin_border = Border(
        left=Side(border_style=BORDER_THIN, color='000000'),
        right=Side(border_style=BORDER_THIN, color='000000'),
        top=Side(border_style=BORDER_THIN, color='000000'),
        bottom=Side(border_style=BORDER_THIN, color='000000')
    )

    no_fill = PatternFill(fill_type=None)
    no_border = Border(
        left=Side(border_style=None, color='FFFFFFFF'),
        right=Side(border_style=None, color='FFFFFFFF'),
        top=Side(border_style=None, color='FFFFFFFF'),
        bottom=Side(border_style=None, color='FFFFFFFF'),
    )
    for column_cells in WorkSheet.rows:
        for cell in column_cells:

            if cell.value or cell.value == 0:
                cell.font = Font(size=11)
                cell.fill = PatternFill("solid", fgColor="f3f3f3")
                cell.border = thin_border


@csrf_exempt
def downloadPPT(request, *args, **kwargs):

    EstimateData = JSONParser().parse(request)

    # return JsonResponse(EstimateData['PlanDetail'], safe=False)
    # Create a presentation instance
    presentation = Presentation()

    # Set slide width and height
    slides_width = Inches(15)
    slides_height = Inches(9)
    presentation.slide_width = slides_width
    presentation.slide_height = slides_height
    # Define slide layouts
    slide = setOverview(presentation, slides_width,
                        slides_height, EstimateData)
    for mediaReportData in EstimateData['Report']:
        slide = setCampaignDetail(
            presentation, slides_width, slides_height, mediaReportData)

    slide = lastSlide(presentation, slides_width, slides_height, EstimateData)

    # Save the presentation
    isExist = os.path.exists('media/estimate/excel/')
    if not isExist:
        os.makedirs('media/estimate/excel/')
    presentation.save('media/estimate/excel/' +
                      EstimateData['PlanDetail']['PlanName']+"-Report.pptx")

    return JsonResponse('/media/estimate/excel/'+EstimateData['PlanDetail']['PlanName']+"-Report.pptx", safe=False)


def setOverview(presentation, slides_width, slides_height, EstimateData):

    # Slide 1: Title slide
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)

    # Add a background image shape
    background = slide.shapes.add_picture(
        "images/estimatepage-bk.png", Inches(0), Inches(0), slides_width, slides_height)
    background.z_order = -1

    title_shape = slide.shapes.title
    title_shape.text = " "

    # Add text box
    width = Inches(3.4)
    height = Inches(1)
    left = Inches(0.7)
    top = Inches(0.2)
    slideHead = slide.shapes.add_textbox(left, top, width, height)
    tf = slideHead.text_frame
    text = tf.add_paragraph()
    text.text = "Campaign Tracker"
    text.font.bold = True
    text.font.size = Pt(40)
    text = tf.add_paragraph()

    text.text = "as on "+datetime.today().strftime('%d-%m-%Y')
    text.font.bold = True
    text.font.size = Pt(20)

    # Set the starting and ending coordinates for the line
    start_x = Inches(0.2)
    start_y = Inches(1.1)
    end_x = Inches(1.2)
    end_y = start_y
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
    line.rotation = 90
    line.line.color.rgb = RGBColor(255, 143, 48)
    line.line.width = Pt(5)

    # Check if the slide has a subtitle
    if len(slide.placeholders) > 1:
        # Remove the subtitle shape
        subtitle = slide.placeholders[1]
        slide.shapes._spTree.remove(subtitle._element)

    # Add an image
    logo_width = Inches(3)
    left = slides_width - logo_width - Inches(0.5)
    top = Inches(0.5)
    slide.shapes.add_picture("images/logo.png", left, top, logo_width)

    width = slides_width
    height = Inches(2)
    left = (slides_width/2)-width/2
    top = (slides_height/2)-height
    slideHead = slide.shapes.add_textbox(left, top, width, height)
    tf = slideHead.text_frame
    text = tf.add_paragraph()
    text.alignment = PP_ALIGN.CENTER
    text.text = "Overview"
    text.font.bold = True
    text.font.size = Pt(40)

    width = Inches(1)
    logo_width = Inches(1)
    height = Inches(1)
    count = EstimateData['UnitCount']
    center = (slides_width/2)
    first_position = 0
    extra_gap = Inches(6)/(count-1)
    top = (slides_height/2)-height
    if count % 2 == 0:
        left_total_gap = extra_gap*((count-1)/2)
        first_position = center-(width*(count/2))-left_total_gap
    else:
        left_total_gap = extra_gap*((count-1)/2)
        first_position = center-(width/2)-(width*((count-1)/2))-left_total_gap
    for i in range(count):
        left = first_position+(width*i)+(extra_gap*i)
        try:
            slide.shapes.add_picture(
                str(EstimateData['Overview'][i]['CampaignUnitIconSRC']), left, top+Inches(0.5), logo_width)
        except:
            slide.shapes.add_picture(
                "images/icon.png", left, top+Inches(0.5), logo_width)
        slideHead = slide.shapes.add_textbox(
            left, top+Inches(1.2), width, height)
        tf = slideHead.text_frame
        text = tf.add_paragraph()
        text.alignment = PP_ALIGN.CENTER
        text.text = EstimateData['Overview'][i]['CampaignUnit']+"\nTotal Served - " + \
            str(EstimateData['Overview'][i]['TotalServed']) + \
            "\nRemaining - "+str(EstimateData['Overview'][i]['Remaining'])+""
        # text.font.bold = True
        text.font.size = Pt(17)

    # bottom left plan details
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    left = Inches(0.5)
    top = slides_height-Inches(2)
    width = Inches(4)
    height = Inches(2)
    topdate_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = topdate_box.text_frame
    # paragraph = text_frame.paragraphs[0]
    # paragraph.line_spacing = Pt(20)
    pladetailHead = text_frame.add_paragraph()
    pladetailHead.text = "Plan Details"
    pladetailHead.font.size = Pt(25)
    pladetailHead.font.bold = True
    planName = text_frame.add_paragraph()
    planName.text = "Plan Name: "+EstimateData['PlanDetail']['PlanName']
    planName.font.size = Pt(18)


def setCampaignDetail(presentation, slides_width, slides_height, MediaReportData):
    detailsimage_slide_layout = presentation.slide_layouts[2]
    slide = presentation.slides.add_slide(detailsimage_slide_layout)

    # Remove main title
    title_shape = slide.shapes.title
    title_shape.text = " "

    # Check if the slide has a subtitle
    if len(slide.placeholders) > 1:
        # Remove the subtitle shape
        subtitle = slide.placeholders[1]
        slide.shapes._spTree.remove(subtitle._element)

    # item heading
    left = Inches(0.5)
    top = Inches(0.2)
    width = slides_width-Inches(0.5)
    height = Inches(1)
    table_head = slide.shapes.add_textbox(left, top, width, height)
    text_frame = table_head.text_frame
    itemHead = text_frame.add_paragraph()
    itemHead.text = MediaReportData['MediaName']
    # itemHead.text = "CampaignPlanDetail['MediaDetail']['MediaName'] +
    #     " - "+CampaignPlanDetail['MediaAdTypeDetail']['MediaTypeName']"
    itemHead.font.size = Pt(22)
    itemHead.font.bold = True
    itemCountry = text_frame.add_paragraph()
    itemCountry.text = MediaReportData['Country']
    # itemCountry.text = CampaignPlanDetail['MediaDetail']['CountryEvent']
    itemCountry.font.size = Pt(16)
    itemType = text_frame.add_paragraph()
    itemType.text = MediaReportData['VerticalName']
    # itemType.text = CampaignPlanDetail['VerticalName']
    itemType.font.size = Pt(15)

    logo_width = Inches(3)
    left = slides_width - logo_width - Inches(0.5)
    top = Inches(0.2)
    try:
        slide.shapes.add_picture(
            MediaReportData['MediaImage'], left, top+Inches(0.5), logo_width)
    except:
        print("")
        # slide.shapes.add_picture(
        #     "images/icon.png", left, top+Inches(0.5), logo_width)
    # Set the starting and ending coordinates for the line
    start_x = Inches(0.5)
    start_y = Inches(1.5)
    end_x = slides_width - logo_width - Inches(1)
    end_y = start_y
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
    line.line.color.rgb = RGBColor(255, 143, 48)
    line.line.width = Pt(2)

    count = len(MediaReportData['CampaignReport'])
    left = -Inches(1.5)
    row_count=3
    top = Inches(3)
    width = (slides_width/row_count)
    height = Inches(1.5)
    cell = 0
    for i in range(count):
        if i % row_count == 0 and i > 0:
            left =-Inches(1.5)
            top = top+height+Inches(1)
            cell = 0

        chart_data = ChartData()
        chart_data_categories = ['Total Served : \n'+str(MediaReportData['CampaignReport'][i]['ServedData'][0]), 'Total Served : \n'+str(MediaReportData['CampaignReport'][i]['ServedData'][1])]
        chart_data_add_series = MediaReportData['CampaignReport'][i]['series']
        
        if chart_data_categories:
            chart_data.categories = chart_data_categories
            chart_data.add_series('', (chart_data_add_series))
            x, y, cx, cy = left + width*cell + \
                (Inches(0.1)*cell), top-Inches(0.1), width + \
                Inches(1), height+Inches(1)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
            ).chart
            
            chart.chart_title.text_frame.text=MediaReportData['CampaignReport'][i]['CampaignUnit']
            chart.has_legend = True
            chart.font.size = Pt(15)
            cell = cell+1


def lastSlide(presentation, slides_width, slides_height, EstimateData):

    # Slide 8: Total table slide
    itemdetail_slide_layout = presentation.slide_layouts[2]
    slide = presentation.slides.add_slide(itemdetail_slide_layout)

    # Remove main title
    title_shape = slide.shapes.title
    title_shape.text = " "

    # Check if the slide has a subtitle
    if len(slide.placeholders) > 1:
        # Remove the subtitle shape
        subtitle = slide.placeholders[1]
        slide.shapes._spTree.remove(subtitle._element)

    background_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        slides_width,
        Inches(4.5),
    )
    background_shape.fill.solid()
    background_shape.fill.fore_color.rgb = RGBColor(
        239, 239, 239)  # Set background color
    background_shape.shadow.inherit = False
    background_shape.shadow.visible = True

    # Add a heading
    heading_text_box = slide.shapes.add_textbox(
        left = Inches(0),
        top=Inches(1.5),
        width=slides_width,
        height=Inches(1.5),
    )
    heading_text_frame = heading_text_box.text_frame
    heading_paragraph = heading_text_frame.add_paragraph()
    heading_paragraph.text = "Thank You"
    heading_paragraph.font.bold = True
    heading_paragraph.font.size = Pt(50)
    heading_paragraph.alignment = PP_ALIGN.CENTER

    left = Inches(0.5)
    top = Inches(4.5)
    width = Inches(14)
    height = Inches(1)
    table_head = slide.shapes.add_textbox(left, top, width, height)
    text_frame = table_head.text_frame
    itemHead = text_frame.add_paragraph()
    itemHead.text = " Masscom Global FZE, Office No: 220, One UAQ Building, UAQ Free Trade Zone, Umm Al Quwain, United Arab Emirates. "
    itemHead.font.size = Pt(18)
    itemHead.alignment = PP_ALIGN.CENTER
    itemHead.word_wrap = True

    # Add an image
    image_path = "images/mail.png"
    image_left = (slides_width/2)-Inches(5)/2 - Inches(0.4)
    image_top = Inches(5.4)
    image_width = Inches(0.4)
    image_height = Inches(0.4)
    slide.shapes.add_picture(image_path, image_left,
                             image_top, image_width, image_height)

    left = (slides_width/2)-Inches(5)/2
    top = Inches(5.1)
    width = Inches(2.5)
    height = Inches(1)
    table_head = slide.shapes.add_textbox(left, top, width, height)
    text_frame = table_head.text_frame
    itemHead = text_frame.add_paragraph()
    itemHead.text = "cs@globalmediakit.com"
    itemHead.font.size = Pt(18)
    itemHead.alignment = PP_ALIGN.CENTER
    itemHead.word_wrap = True

    # Add an image
    image_path = "images/phone.png"
    image_left = ((slides_width/2)-Inches(5)/2) + Inches(2.9)
    image_top = Inches(5.4)
    image_width = Inches(0.4)
    image_height = Inches(0.4)
    slide.shapes.add_picture(image_path, image_left,
                             image_top, image_width, image_height)

    left = ((slides_width/2)-Inches(5)/2) + Inches(2.9)
    top = Inches(5.1)
    width = Inches(2.5)
    height = Inches(1)
    table_head = slide.shapes.add_textbox(left, top, width, height)
    text_frame = table_head.text_frame
    itemHead = text_frame.add_paragraph()
    itemHead.text = "+971-43827760"
    itemHead.font.size = Pt(18)
    itemHead.alignment = PP_ALIGN.CENTER
    itemHead.word_wrap = True

    left = Inches(0.5)
    top = Inches(5.9)
    width = Inches(14)
    height = Inches(1)
    table_head = slide.shapes.add_textbox(left, top, width, height)
    text_frame = table_head.text_frame
    itemHead = text_frame.add_paragraph()
    itemHead.text = "Memberships"
    itemHead.font.size = Pt(20)
    itemHead.alignment = PP_ALIGN.CENTER
    itemHead.word_wrap = True

    background_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(6.6),
        Inches(14),
        Inches(2),
    )
    background_shape.fill.solid()
    background_shape.fill.fore_color.rgb = RGBColor(
        239, 239, 239)  # Set background color
    background_shape.shadow.inherit = False
    background_shape.shadow.visible = True

    # Add an image
    image_path = "images/h1.png"
    image_left = Inches(0.75)
    image_top = Inches(6.9)
    image_width = Inches(2.3)
    image_height = Inches(1.4)
    slide.shapes.add_picture(image_path, image_left,
                             image_top, image_width, image_height)

    # Add an image
    image_path = "images/h2.png"
    image_left = Inches(3.55)
    image_top = Inches(6.9)
    image_width = Inches(2.3)
    image_height = Inches(1.4)
    slide.shapes.add_picture(image_path, image_left,
                             image_top, image_width, image_height)

    # Add an image
    image_path = "images/h3.png"
    image_left = Inches(6.35)
    image_top = Inches(6.9)
    image_width = Inches(2.3)
    image_height = Inches(1.4)
    slide.shapes.add_picture(image_path, image_left,
                             image_top, image_width, image_height)

    # Add an image
    image_path = "images/h4.png"
    image_left = Inches(9.15)
    image_top = Inches(6.9)
    image_width = Inches(2.3)
    image_height = Inches(1.4)
    slide.shapes.add_picture(image_path, image_left,
                             image_top, image_width, image_height)

    # Add an image
    image_path = "images/h5.png"
    image_left = Inches(11.95)
    image_top = Inches(6.9)
    image_width = Inches(2.3)
    image_height = Inches(1.4)
    slide.shapes.add_picture(image_path, image_left,
                             image_top, image_width, image_height)
    return slide
