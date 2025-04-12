import os
from django.http import JsonResponse
from django.http import JsonResponse
from rest_framework.parsers import JSONParser
# Create your views here.
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from pptx.enum.shapes import MSO_CONNECTOR
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
import requests
from io import BytesIO
from django.views.decorators.csrf import csrf_exempt

from app_vertical_media.models import MediaAdTypeImage
from app_vertical_media.serializers import UploadMediaAdTypeImageSerializer


@csrf_exempt
def estimatePPTDownload(request, *args, **kwargs):

    EstimateData = JSONParser().parse(request)

    # return JsonResponse(EstimateData['PlanDetail'], safe=False)
    # Create a presentation instance
    presentation = Presentation()

    # Set slide width and height
    slides_width = Inches(15)
    slides_height = Inches(9)
    presentation.slide_width = slides_width
    presentation.slide_height = slides_height

    # Define slide layouts
    title_slide_layout = presentation.slide_layouts[0]

    # Slide 1: Title slide
    slide = presentation.slides.add_slide(title_slide_layout)

    # Add a background image shape
    background = slide.shapes.add_picture(
        "images/estimatepage-bk.png", Inches(0), Inches(0), slides_width, slides_height)
    background.z_order = -1

    title_shape = slide.shapes.title
    title_shape.text = " "

    # Add text box
    width = Inches(3.4)
    height = Inches(1)
    left = (slides_width/2)-width/2
    top = (slides_height/2)-height
    slideHead = slide.shapes.add_textbox(left, top, width, height)
    tf = slideHead.text_frame
    text = tf.add_paragraph()
    text.text = "Media Plan"
    text.font.bold = True
    text.font.size = Pt(50)

    # Set the starting and ending coordinates for the line
    start_x = Inches(5.5)
    start_y = Inches(4.28)
    end_x = Inches(6)
    end_y = start_y
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
    line.rotation = 90
    line.line.color.rgb = RGBColor(255, 143, 48)
    line.line.width = Pt(5)

    # Check if the slide has a subtitle
    if len(slide.placeholders) > 1:
        # Remove the subtitle shape
        subtitle = slide.placeholders[1]
        slide.shapes._spTree.remove(subtitle._element)

    # top left dates
    left = Inches(0.5)
    top = Inches(0.2)
    width = Inches(4)
    height = Inches(2)
    topdate_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = topdate_box.text_frame
    createDate = text_frame.add_paragraph()
    createDate.text = "Plan Created on: " + \
        EstimateData['PlanDetail']['PlanCreatedDate']
    createDate.font.size = Pt(18)
    expiryDate = text_frame.add_paragraph()
    expiryDate.text = "Expiry date: " + \
        EstimateData['PlanDetail']['PlanExpiryDate']
    expiryDate.font.size = Pt(18)

    # Add an image
    logo_width = Inches(3)
    left = slides_width - logo_width - Inches(0.5)
    top = Inches(0.5)
    slide.shapes.add_picture("images/logo.png", left, top, logo_width)

    # bottom left plan details
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    left = Inches(0.5)
    top = slides_height-Inches(2)
    width = Inches(4)
    height = Inches(2)
    topdate_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = topdate_box.text_frame
    # paragraph = text_frame.paragraphs[0]
    # paragraph.line_spacing = Pt(20)
    pladetailHead = text_frame.add_paragraph()
    pladetailHead.text = "Plan Details"
    pladetailHead.font.size = Pt(25)
    pladetailHead.font.bold = True
    planName = text_frame.add_paragraph()
    planName.text = "Plan Name: "+EstimateData['PlanDetail']['PlanName']
    planName.font.size = Pt(18)
    planId = text_frame.add_paragraph()
    planId.text = "Plan Id: "+EstimateData['PlanDetail']['PlanId']
    planId.font.size = Pt(18)
    startDate = text_frame.add_paragraph()
    startDate.text = "Tentative Start Date: " + \
        EstimateData['PlanDetail']['StartDate']
    startDate.font.size = Pt(18)

    # bottom right details
    left = slides_width-Inches(4)
    top = slides_height-Inches(2)
    width = Inches(4)
    height = Inches(2)
    topdate_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = topdate_box.text_frame
    pladetailHead = text_frame.add_paragraph()
    pladetailHead.text = "Plan created by"
    pladetailHead.font.size = Pt(25)
    pladetailHead.font.bold = True
    planName = text_frame.add_paragraph()
    planName.text = "Name: "+EstimateData['PlanDetail']['PlancreatedName']
    planName.font.size = Pt(18)
    startDate = text_frame.add_paragraph()
    startDate.text = "Company Name: "+EstimateData['PlanDetail']['CompanyName']
    startDate.font.size = Pt(18)
    startDate = text_frame.add_paragraph()
    startDate.text = "Location: "+EstimateData['PlanDetail']['Location']
    startDate.font.size = Pt(18)

    # Slide 2: Content slide with a table
    content_slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(content_slide_layout)

    # Remove main title
    title_shape = slide.shapes.title
    title_shape.text = " "

    # Check if the slide has a subtitle
    if len(slide.placeholders) > 1:
        # Remove the subtitle shape
        subtitle = slide.placeholders[1]
        slide.shapes._spTree.remove(subtitle._element)

    # top heading
    left = Inches(0.5)
    top = Inches(0.2)
    width = Inches(4)
    height = Inches(1)
    slide_head = slide.shapes.add_textbox(left, top, width, height)
    text_frame = slide_head.text_frame
    heading = text_frame.add_paragraph()
    heading.text = "Plan summary"
    heading.font.size = Pt(30)

    chart_data = ChartData()
    chart_data_categories = []
    chart_data_add_series = []
    for BudgetAllocationArray in EstimateData['BudgetAllocationArray']:
        if BudgetAllocationArray['Percentage'] == 'NaN':
            BudgetAllocationArray['Percentage']='0.00'
        if BudgetAllocationArray['Percentage'] != '0' and BudgetAllocationArray['Percentage'] != '0.00':
            chart_data_categories.append(BudgetAllocationArray['Name'])
            chart_data_add_series.append(
                float(BudgetAllocationArray['Percentage']))
            
    if chart_data_categories:
        # Define the position and size of the background box
        left = Inches(0.5)
        top = Inches(1.3)
        width = (slides_width/2)-Inches(0.5)
        height = Inches(7)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 240, 240)

        # Set the line color and width of the background box (optional)
        line = shape.line
        line.color.rgb = RGBColor(240, 240, 240)
        line.width = Pt(1)
        shape.z_order = 0

   
        chart_data.categories = chart_data_categories
        chart_data.add_series('', (chart_data_add_series))

        x, y, cx, cy = Inches(.5), Inches(1.7), Inches(7), Inches(6)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = False

        # set labels to contain category and value
        for i in range(len(chart_data.categories)):
            point = chart.series[0].points[i]
            point.data_label.text_frame.text = "{}".format(
                chart_data.categories[i].label)
            point.data_label.position = XL_LABEL_POSITION.OUTSIDE_END
            point.data_label.text_frame.word_wrap = True
            point.data_label.text_frame.paragraphs[0].font.size = Pt(5)
            for run in point.data_label.text_frame.paragraphs[0].runs:
                run.font.size = Pt(10)
        

    # table heading
    if  chart_data_categories:
        left = (slides_width/2)+Inches(0.4)
    else:
        left = Inches(0.5)
    top = Inches(0.9)
    width = Inches(4)
    height = Inches(1)
    table_head = slide.shapes.add_textbox(left, top, width, height)
    text_frame = table_head.text_frame
    tableHeading = text_frame.add_paragraph()
    tableHeading.text = "Budget Allocation"
    tableHeading.font.size = Pt(20)

    # Define table dimensions and position
    if  chart_data_categories:
        left = (slides_width/2)+Inches(0.5)
        width = int((slides_width/2)-Inches(1))
    else:
        left = Inches(0.5)
        width = int((slides_width)-Inches(1))
    rows = 9
    cols = 3
    top = Inches(1.7)
    height = Inches(2.5)

    # Add table to the slide
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set table style
    table.style = "Table Grid"

    # Populate table cells
    table.cell(0, 0).text = "Media"
    table.cell(0, 1).text = "Net Investment(USD)"
    table.cell(0, 2).text = "Budget Percentage"
    cellindex = 0
    for BudgetAllocationArray in EstimateData['BudgetAllocationArray']:
        cellindex = cellindex+1
        table.cell(cellindex, 0).text = BudgetAllocationArray['Name']
        table.cell(cellindex, 1).text = str(BudgetAllocationArray['Amount'])
        table.cell(cellindex, 2).text = BudgetAllocationArray['Percentage']

    table.cell(8, 0).text = "Total investment"
    table.cell(8, 1).text = EstimateData['PlanDetail']['Budget']
    table.cell(8, 2).text = "100"

    column_padding = Inches(0.2)
    # Customize table cell formatting
    for row in range(rows):
        for col in range(cols):
            cell = table.cell(row, col)
            if col == 0:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            else:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

            cell.text_frame.paragraphs[0].font.size = Pt(16)
            cell.margin_left = column_padding
            cell.margin_right = column_padding
            cell.margin_top = column_padding
            cell.margin_bottom = column_padding
            paragraph = text_frame.paragraphs[0]
            paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    header_row = table.rows[0]
    for cell in header_row.cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(220, 220, 220)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    row_colors = [RGBColor(240, 240, 240),
                  RGBColor(240, 240, 240)]
    for i, row in enumerate(table.rows):
        color_index = i % len(row_colors)
        for cell in row.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_colors[color_index]

# ---------------------------------------------------------------------------------

    for BudgetAllocationArray in EstimateData['BudgetAllocationArray']:
        if str(BudgetAllocationArray['Amount'])!='0.00' and str(BudgetAllocationArray['Percentage'])!='0.00' and BudgetAllocationArray['AdvantageImage']:
            # Slide 2.1: Advantage slide
            detailshead_slide_layout = presentation.slide_layouts[2]
            slide = presentation.slides.add_slide(detailshead_slide_layout)

            # Remove main title
            title_shape = slide.shapes.title
            title_shape.text = " "

            # Check if the slide has a subtitle
            if len(slide.placeholders) > 1:
                # Remove the subtitle shape
                subtitle = slide.placeholders[1]
                slide.shapes._spTree.remove(subtitle._element)

            # Add text box
            
            
            left = Inches(0.5)
            top = Inches(0.2)
            width = slides_width-Inches(0.5)
            height = Inches(1)
            slideHead = slide.shapes.add_textbox(left, top, width, height)
            tf = slideHead.text_frame
            text = tf.add_paragraph()
            text.text = "Advantage of "+BudgetAllocationArray['Name']
            text.font.bold = True
            text.font.size = Pt(30)
            start_x = Inches(0.5)
            start_y = Inches(1.2)
            end_x = slides_width-Inches(0.5)
            end_y = start_y
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
            line.line.color.rgb = RGBColor(255, 143, 48)
            line.line.width = Pt(2)

            left = (slides_width-Inches(12))/2
            top = Inches(1.5)
            width = Inches(12)
            height = Inches(7)
            response = requests.get(BudgetAllocationArray['AdvantageImage'])
            image_data = BytesIO(response.content)
            slide.shapes.add_picture(image_data, left, top, width, height)

# ---------------------------------------------------------------------------------


    # Slide 3: Details title slide
    detailshead_slide_layout = presentation.slide_layouts[2]
    slide = presentation.slides.add_slide(detailshead_slide_layout)

    # Remove main title
    title_shape = slide.shapes.title
    title_shape.text = " "

    # Check if the slide has a subtitle
    if len(slide.placeholders) > 1:
        # Remove the subtitle shape
        subtitle = slide.placeholders[1]
        slide.shapes._spTree.remove(subtitle._element)

    # Add text box
    width = Inches(5)
    height = Inches(1)
    left = (slides_width/2)-width/2
    top = (slides_height/2)-height
    slideHead = slide.shapes.add_textbox(left, top, width, height)
    tf = slideHead.text_frame
    text = tf.add_paragraph()
    text.text = "SELECTED OPTIONS"
    text.font.bold = True
    text.font.size = Pt(50)

    # Set the starting and ending coordinates for the line
    start_x = Inches(4.5)
    start_y = Inches(4.28)
    end_x = Inches(5)
    end_y = start_y
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
    line.rotation = 90
    line.line.color.rgb = RGBColor(255, 143, 48)
    line.line.width = Pt(5)
    for ListByCountry in EstimateData['CampaignDetailFilterByCountry']:
        # Slide 5: Details title slide
        countrycover_slide_layout = presentation.slide_layouts[2]
        slide = presentation.slides.add_slide(countrycover_slide_layout)
        # Remove main title
        title_shape = slide.shapes.title
        title_shape.text = " "

        # Check if the slide has a subtitle
        if len(slide.placeholders) > 1:
            # Remove the subtitle shape
            subtitle = slide.placeholders[1]
            slide.shapes._spTree.remove(subtitle._element)
         # Add text box
        width = slides_width-Inches(1)
        height = Inches(2)
        left = Inches(.5)
        top = Inches(3)
        slideHead = slide.shapes.add_textbox(left, top, width, height)
        tf = slideHead.text_frame
        tf.word_wrap = True
        text = tf.add_paragraph()
        text.text = ListByCountry['CountryName']
        text.font.bold = True
        text.font.size = Pt(40)
        text.alignment = PP_ALIGN.CENTER


        
        width = slides_width
        height = Inches(2)
        left = Inches(0)
        top = Inches(1.6)
        slideHead = slide.shapes.add_textbox(left, top, width, height)
        tf = slideHead.text_frame
        tf.word_wrap = True
        text = tf.add_paragraph()
        text.text = "Countries/Events Covered"
        text.font.bold = True
        text.font.size = Pt(40)
        text.alignment = PP_ALIGN.CENTER

        # Set the starting and ending coordinates for the line
        start_x = Inches(4.9)
        start_y = Inches(2.7)
        end_x = Inches(10.2)
        end_y = start_y
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
        # line.rotation = 90
        line.line.color.rgb = RGBColor(255, 143, 48)
        line.line.width = Pt(5)
        for CampaignPlanDetail in ListByCountry['CampaignDetails']:

            # Slide 4: Details title slide
            detailsimage_slide_layout = presentation.slide_layouts[2]
            slide = presentation.slides.add_slide(detailsimage_slide_layout)

            # Remove main title
            title_shape = slide.shapes.title
            title_shape.text = " "

            # Check if the slide has a subtitle
            if len(slide.placeholders) > 1:
                # Remove the subtitle shape
                subtitle = slide.placeholders[1]
                slide.shapes._spTree.remove(subtitle._element)

            # item heading
            left = Inches(0.5)
            top = Inches(0.2)
            width = slides_width-Inches(0.5)
            height = Inches(1)
            table_head = slide.shapes.add_textbox(left, top, width, height)
            text_frame = table_head.text_frame
            itemHead = text_frame.add_paragraph()
            itemHead.text = CampaignPlanDetail['MediaDetail']['MediaName'] + \
                " - "+CampaignPlanDetail['MediaAdTypeDetail']['MediaTypeName']
            itemHead.font.size = Pt(22)
            itemCountry = text_frame.add_paragraph()
            itemCountry.text = CampaignPlanDetail['MediaDetail']['CountryEvent']
            itemCountry.font.size = Pt(18)
            itemType = text_frame.add_paragraph()
            itemType.text = CampaignPlanDetail['VerticalName']
            itemType.font.size = Pt(18)

            # Set the starting and ending coordinates for the line
            start_x = Inches(0.5)
            start_y = Inches(2)
            end_x = slides_width-Inches(0.5)
            end_y = start_y
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
            line.line.color.rgb = RGBColor(255, 143, 48)
            line.line.width = Pt(2)

            mediaAdTypeImage=MediaAdTypeImage.objects.filter(
            MediaAdTypeId=CampaignPlanDetail['MediaAdTypeDetail']['MediaAdTypeId']) 
            mediaSerializer = UploadMediaAdTypeImageSerializer(
                mediaAdTypeImage, many=True)
            adUnitImage=mediaSerializer.data
            i=0
            for AdunitImage in adUnitImage:
                # Add image1
                if i==0 :
                    left = (slides_width-Inches(10))/2
                    top = Inches(2.5)
                    width = Inches(10)
                    height = Inches(6)
                    
                    url = EstimateData['PlanDetail']['ImagePath'] + \
                        AdunitImage['MediaAdTypeImage']
                    response = requests.get(url)
                    image_data = BytesIO(response.content)

                    try:
                        slide.shapes.add_picture(image_data, left, top, width, height)

                    except:
                        response = requests.get(
                            "https://globalmediakit.com/assets/images/image_place_holder.png")
                        image_data = BytesIO(response.content)
                        slide.shapes.add_picture(image_data, left, top, width, height)
                else :
                    # Add image 2
                    slide = presentation.slides.add_slide(detailsimage_slide_layout)
                    # Remove main title
                    title_shape = slide.shapes.title
                    title_shape.text = " "

                    # Check if the slide has a subtitle
                    if len(slide.placeholders) > 1:
                        # Remove the subtitle shape
                        subtitle = slide.placeholders[1]
                        slide.shapes._spTree.remove(subtitle._element)
                    left = Inches(0)
                    top = Inches(0)
                    width = slides_width
                    height = slides_height

                    url = EstimateData['PlanDetail']['ImagePath'] + \
                        AdunitImage['MediaAdTypeImage']
                    response = requests.get(url)
                    image_data = BytesIO(response.content)

                    try:
                        slide.shapes.add_picture(image_data, left, top, width, height)

                    except:
                        response = requests.get(
                            "https://globalmediakit.com/assets/images/image_place_holder.png")
                        image_data = BytesIO(response.content)
                        slide.shapes.add_picture(image_data, left, top, width, height)

                i=i+1
                

            # Slide 5: Details title slide
            itemdetail_slide_layout = presentation.slide_layouts[2]
            slide = presentation.slides.add_slide(itemdetail_slide_layout)

            # Remove main title
            title_shape = slide.shapes.title
            title_shape.text = " "

            # Check if the slide has a subtitle
            if len(slide.placeholders) > 1:
                # Remove the subtitle shape
                subtitle = slide.placeholders[1]
                slide.shapes._spTree.remove(subtitle._element)

            # item heading
            left = Inches(0.5)
            top = Inches(0.2)
            width = slides_width-Inches(1)
            height = Inches(1)
            table_head = slide.shapes.add_textbox(left, top, width, height)
            text_frame = table_head.text_frame
            itemHead = text_frame.add_paragraph()
            itemHead.text = CampaignPlanDetail['MediaDetail']['MediaName'] + \
                " - "+CampaignPlanDetail['MediaAdTypeDetail']['MediaTypeName']
            itemHead.font.size = Pt(22)
            text_frame.word_wrap = True
            itemCountry = text_frame.add_paragraph()
            itemCountry.text = CampaignPlanDetail['MediaDetail']['CountryEvent']
            itemCountry.font.size = Pt(18)
            itemType = text_frame.add_paragraph()
            itemType.text = CampaignPlanDetail['VerticalName']
            itemType.font.size = Pt(18)

            # Add a text box
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(3)
            height = Inches(1)
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame

            # Add the text to the text box
            text = "More Details"
            p = text_frame.add_paragraph()
            run = p.add_run()
            run.text = text

            # Add a hyperlink to the text
            hyperlink = run.hyperlink
            if CampaignPlanDetail['MediaAdUnitLink']:
                hyperlink.address = CampaignPlanDetail['MediaAdUnitLink']
            else:
                hyperlink.address = CampaignPlanDetail['MediaLink']

            # Set the starting and ending coordinates for the line
            start_x = Inches(0.5)
            start_y = Inches(2.4)
            end_x = Inches(14.5)
            end_y = start_y
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
            line.line.color.rgb = RGBColor(255, 143, 48)
            line.line.width = Pt(1)

            # rate head
            left = Inches(0.4)
            top = Inches(2.5)
            width = slides_width - Inches(1)
            height = Inches(1)
            table_head = slide.shapes.add_textbox(left, top, width, height)
            text_frame = table_head.text_frame
            text_frame.word_wrap = True
            itemDetails = text_frame.add_paragraph()
            itemDetails.text = "Rates and details"
            itemDetails.font.size = Pt(22)
            itemDetails.alignment = PP_ALIGN.LEFT

            # Define table dimensions and position
            rows = 6
            cols = 1
            left = Inches(0.5)
            top = Inches(3.3)
            width = int(slides_width/2 - Inches(0.5))
            height = Inches(2.5)
            # Add table to the slide
            table = slide.shapes.add_table(
                rows, cols, left, top, width, height).table
            # Set table style
            table.style = "Table Grid"
            table.first_row = False
            dataLeft = []
            dataRight = []
            if CampaignPlanDetail['MediaAdTypeDetail']['MediaAdType'] == 'Ad Unit' and CampaignPlanDetail['MediaDetail']['VerticalName'] != 'Media Service':
                if CampaignPlanDetail['MediaDetail']['VerticalName'] == 'Out of Home':
                    dataLeft = [
                        ["Start Date: "+CampaignPlanDetail['StartDate']],
                        ["Number of Screens: " +
                            CampaignPlanDetail['MediaAdTypeDetail']['NoOfScreen']],
                        ["Lighting type: "+CampaignPlanDetail['MediaAdTypeDetail']
                            ['PositionName']],
                        ["Size: " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['OutdoorSize'])],
                        ["Single Rotation Time in Seconds: " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['SingleRotationTime'])],
                        ["Duration of creative in seconds: " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['DurationOfCreative'])]
                        # Add more rows as needed
                    ]
    
                    if CampaignPlanDetail['MediaAdTypeDetail']['RentalUsd'] !='0.00':
                        dataRight = [
                            ["Number of production: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfProduction'])],
                            ["One time production cost: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['ProductionCostUsd'])],
                            ["Rental Cost/Month: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['RentalUsd'])],
                            ["Discounted rate/Month (Rental): USD " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['DiscountUsd'])],
                            ["Total number of months required: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfMonths'])],
                            ["Sub Total (Rental+ One Time Production): USD " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['SubTotal'])]
                            # Add more rows as needed
                        ]
                    else :
                        dataRight = [
                            ["Number of production: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfProduction'])],
                            ["One time production cost: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['ProductionCostUsd'])],
                            ["Rental Cost/Month: " +
                                "Final Rates To Be Confirmed"],
                            ["Discounted rate/Month (Rental): " +
                            "Final Rates To Be Confirmed"],
                            ["Total number of months required: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfMonths'])],
                            ["Sub Total (Rental+ One Time Production): " +
                            "Final Rates To Be Confirmed"]
                            # Add more rows as needed
                        ]
                if CampaignPlanDetail['MediaDetail']['VerticalName'] == 'Television' or CampaignPlanDetail['MediaDetail']['VerticalName'] == 'Radio':
                    dataLeft = [
                        ["Time band: "+CampaignPlanDetail['MediaAdTypeDetail']['PositionName']],
                        ["Spot Duration (seconds): " +
                            CampaignPlanDetail['MediaAdTypeDetail']['SizeName']],
                        ["Time: "+CampaignPlanDetail['MediaAdTypeDetail']['ProgramStartTime'] +
                            " - "+CampaignPlanDetail['MediaAdTypeDetail']['ProgramEndTime']],
                        # ["Rate/spot: "+str(CampaignPlanDetail['MediaAdTypeDetail']['RentalUsd'])],
                        # ["Number of production: " +
                        #     str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfProduction'])],
                        # ["One time production cost: " +
                        #     str(CampaignPlanDetail['MediaAdTypeDetail']['ProductionCostUsd'])]
                        # Add more rows as needed
                    ]

                    
                    
                    if CampaignPlanDetail['MediaAdTypeDetail']['RentalUsd'] !='0.00':
                        dataRight = [
                            ["Number of production: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfProduction'])],
                            ["One time production cost: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['ProductionCostUsd'])],
                            ["Rate/spot: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['RentalUsd'])],
                            ["Discounted rate: USD " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['DiscountUsd'])],
                            ["Number of spots required: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfMonths'])],
                            ["Sub Total: USD " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['SubTotal'])]
                            # Add more rows as needed
                        ]
                    else :
                        dataRight = [
                            ["Number of production: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfProduction'])],
                            ["One time production cost: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['ProductionCostUsd'])],
                            ["Rate/spot: " +
                                "Final Rates To Be Confirmed"],
                            ["Discounted rate: " +
                            "Final Rates To Be Confirmed"],
                            ["Number of spots required: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfMonths'])],
                            ["Sub Total: " +
                            "Final Rates To Be Confirmed"]
                            # Add more rows as needed
                        ]

                if CampaignPlanDetail['MediaDetail']['VerticalName'] == 'Newspaper' or CampaignPlanDetail['MediaDetail']['VerticalName'] == 'Magazine':
                    dataLeft = [
                        ["Frequency: "+CampaignPlanDetail['MediaDetail']['CategoryName']],
                        ["Size: " +
                            CampaignPlanDetail['MediaAdTypeDetail']['CreativeSize']],
                        [CampaignPlanDetail['MediaAdTypeDetail']['PositionName'] +
                            ": "+CampaignPlanDetail['MediaAdTypeDetail']['SizeName']],
                        # ["Size: "+str(CampaignPlanDetail['MediaAdTypeDetail']['OutdoorSize'])],
                        # ["Single Rotation Time in Seconds: " +
                        #     str(CampaignPlanDetail['MediaAdTypeDetail']['SingleRotationTime'])],
                        # ["Duration of creative in seconds: " +
                        #     str(CampaignPlanDetail['MediaAdTypeDetail']['DurationOfCreative'])]
                        # Add more rows as needed
                    ]

                    
                    if CampaignPlanDetail['MediaAdTypeDetail']['RentalUsd'] !='0.00':
                        dataRight = [
                            ["Number of production: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfProduction'])],
                            ["One time production cost: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['ProductionCostUsd'])],
                            ["Rate/Insertion: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['RentalUsd'])],
                            ["Discounted rate: USD " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['DiscountUsd'])],
                            ["Number of insertions required: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfMonths'])],
                            ["Sub Total (Rental+ One Time Production): USD " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['SubTotal'])]
                            # Add more rows as needed
                        ]
                    else :
                        dataRight = [
                            ["Number of production: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfProduction'])],
                            ["One time production cost: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['ProductionCostUsd'])],
                            ["Rate/Insertion: " +
                                "Final Rates To Be Confirmed"],
                            ["Discounted rate: " +
                            "Final Rates To Be Confirmed"],
                            ["Number of insertions required: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfMonths'])],
                            ["Sub Total: " +
                            "Final Rates To Be Confirmed"]
                            # Add more rows as needed
                        ]

                if CampaignPlanDetail['MediaDetail']['VerticalName'] == 'Digital':
                    dataLeft = [
                        ["Category: "+CampaignPlanDetail['MediaDetail']['CategoryName']],
                        [CampaignPlanDetail['MediaAdTypeDetail']['PositionName'] +
                            ": "+CampaignPlanDetail['MediaAdTypeDetail']['SizeName']],
                        ["Size: " + CampaignPlanDetail['MediaAdTypeDetail']['SizeName']],
                        # ["Size: "+str(CampaignPlanDetail['MediaAdTypeDetail']['OutdoorSize'])],
                        # ["Single Rotation Time in Seconds: " +
                        #     str(CampaignPlanDetail['MediaAdTypeDetail']['SingleRotationTime'])],
                        # ["Duration of creative in seconds: " +
                        #     str(CampaignPlanDetail['MediaAdTypeDetail']['DurationOfCreative'])]
                        # Add more rows as needed
                    ]

                    
                    if CampaignPlanDetail['MediaAdTypeDetail']['RentalUsd'] !='0.00':
                        dataRight = [
                            ["Number of production: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfProduction'])],
                            ["One time production cost: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['ProductionCostUsd'])],
                            ["Cost per Thousand Impressions: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['RentalUsd'])],
                            ["Discounted rate: USD " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['DiscountUsd'])],
                            ["Number of impressions required: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfMonths'])],
                            ["Sub Total : USD " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['SubTotal'])]
                            # Add more rows as needed
                        ]
                    else :
                        dataRight = [
                            ["Number of production: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfProduction'])],
                            ["One time production cost: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['ProductionCostUsd'])],
                            ["Cost per Thousand Impressions: " +
                                "Final Rates To Be Confirmed"],
                            ["Discounted rate: " +
                            "Final Rates To Be Confirmed"],
                            ["Number of impressions required: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfMonths'])],
                            ["Sub Total: " +
                            "Final Rates To Be Confirmed"]
                            # Add more rows as needed
                        ]
            if CampaignPlanDetail['MediaAdTypeDetail']['MediaAdType'] == 'Package' or CampaignPlanDetail['MediaDetail']['VerticalName'] == 'Media Service':
                if CampaignPlanDetail['MediaDetail']['VerticalName'] == 'Out of Home':
                    dataLeft = [
                        ["Number of Screens: "+CampaignPlanDetail['MediaAdTypeDetail']['NoOfScreen']],
                        
                        # Add more rows as needed
                    ]
                    if CampaignPlanDetail['MediaAdTypeDetail']['RentalUsd'] !='0.00':
                        dataRight = [
                            ["Number of production: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfProduction'])],
                            ["One time production cost: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['ProductionCostUsd'])],
                            ["Package Rental Rate: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['RentalUsd'])],
                            ["Discounted rate/Month (Rental): USD " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['DiscountUsd'])],
                            ["Sub Total (Rental+ One Time Production): USD " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['SubTotal'])]
                            # Add more rows as needed
                        ]
                    else :
                        dataRight = [
                            ["Number of production: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfProduction'])],
                            ["One time production cost: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['ProductionCostUsd'])],
                            ["Package Rental Rate: " +
                                "Final Rates To Be Confirmed"],
                            ["Discounted rate: " +
                            "Final Rates To Be Confirmed"],
                            ["Sub Total (Rental+ One Time Production): " +
                            "Final Rates To Be Confirmed"]
                            # Add more rows as needed
                        ]
                if CampaignPlanDetail['MediaDetail']['VerticalName'] == 'Television' or CampaignPlanDetail['MediaDetail']['VerticalName'] == 'Radio':
                    dataLeft = [
                        ["Time: "+CampaignPlanDetail['MediaAdTypeDetail']['ProgramStartTime'] +
                            " - "+CampaignPlanDetail['MediaAdTypeDetail']['ProgramEndTime']],
                        ["Estimated number of spots: " +
                            CampaignPlanDetail['MediaAdTypeDetail']['NumberOfSpots']],
                        # Add more rows as needed
                    ]

                    if CampaignPlanDetail['MediaAdTypeDetail']['RentalUsd'] !='0.00':
                        dataRight = [
                            ["Number of production: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfProduction'])],
                            ["One time production cost: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['ProductionCostUsd'])],
                            ["Package/Sponsorship Rate: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['RentalUsd'])],
                            ["Discounted rate: USD " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['DiscountUsd'])],
                            ["Sub Total: USD " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['SubTotal'])]
                            # Add more rows as needed
                        ]
                    else :
                        dataRight = [
                            ["Number of production: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfProduction'])],
                            ["One time production cost: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['ProductionCostUsd'])],
                            ["Package/Sponsorship Rate: " +
                                "Final Rates To Be Confirmed"],
                            ["Discounted rate: " +
                            "Final Rates To Be Confirmed"],
                            ["Sub Total: " +
                            "Final Rates To Be Confirmed"]
                            # Add more rows as needed
                        ]
                if CampaignPlanDetail['MediaDetail']['VerticalName'] == 'Newspaper' or CampaignPlanDetail['MediaDetail']['VerticalName'] == 'Magazine':
                    dataLeft = [
                        ["Frequency: "+CampaignPlanDetail['MediaDetail']['CategoryName']],
                                            
                        # Add more rows as needed
                    ]

                    
                    if CampaignPlanDetail['MediaAdTypeDetail']['RentalUsd'] !='0.00':
                        dataRight = [
                            ["Number of production: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfProduction'])],
                            ["One time production cost: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['ProductionCostUsd'])],
                            ["Package/Sponsorship Rate: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['RentalUsd'])],
                            ["Discounted rate: USD " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['DiscountUsd'])],
                            ["Sub Total: USD " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['SubTotal'])]
                            # Add more rows as needed
                        ]
                    else :
                        dataRight = [
                            ["Number of production: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfProduction'])],
                            ["One time production cost: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['ProductionCostUsd'])],
                            ["Package/Sponsorship Rate: " +
                                "Final Rates To Be Confirmed"],
                            ["Discounted rate: " +
                            "Final Rates To Be Confirmed"],
                            ["Sub Total: " +
                            "Final Rates To Be Confirmed"]
                            # Add more rows as needed
                        ]
                if CampaignPlanDetail['MediaDetail']['VerticalName'] == 'Digital':
                    dataLeft = [
                        ["Category: "+CampaignPlanDetail['MediaDetail']['CategoryName']],
                        # Add more rows as needed
                    ]

                    
                    if CampaignPlanDetail['MediaAdTypeDetail']['RentalUsd'] !='0.00':
                        dataRight = [
                            ["Number of production: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfProduction'])],
                            ["One time production cost: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['ProductionCostUsd'])],
                            ["Package/Sponsorship Rate: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['RentalUsd'])],
                            ["Discounted rate: USD " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['DiscountUsd'])],
                            ["Sub Total: USD " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['SubTotal'])]
                            # Add more rows as needed
                        ]
                    else :
                        dataRight = [
                            ["Number of production: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfProduction'])],
                            ["One time production cost: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['ProductionCostUsd'])],
                            ["Package/Sponsorship Rate: " +
                                "Final Rates To Be Confirmed"],
                            ["Discounted rate: " +
                            "Final Rates To Be Confirmed"],
                            ["Sub Total: " +
                            "Final Rates To Be Confirmed"]
                            # Add more rows as needed
                        ]
                        
                if CampaignPlanDetail['MediaDetail']['VerticalName'] == 'Media Service':
                    dataLeft = [
                        ["Category: "+CampaignPlanDetail['MediaDetail']['CategoryName']],
                        ["Order Unit: "+CampaignPlanDetail['MediaAdTypeDetail']['OrderUnitName']],
                        ["Ordering Quantity: "+CampaignPlanDetail['MediaAdTypeDetail']['OrderingQuantity']],
                        ["Unit type: "+CampaignPlanDetail['MediaAdTypeDetail']['UnitTypeName']],
                        # Add more rows as needed
                    ]
                    if CampaignPlanDetail['MediaAdTypeDetail']['RentalUsd'] !='0.00':
                
                        dataRight = [
                            ["Number of production: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfProduction'])],
                            ["One time production cost: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['ProductionCostUsd'])],
                            ["Card Rate: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['RentalUsd'])],
                            ["Discounted rate: USD " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['DiscountUsd'])],
                            ["Total number of months required: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfMonths'])],
                            ["Sub Total: USD " +
                            str(CampaignPlanDetail['MediaAdTypeDetail']['SubTotal'])]
                        ]
                    else :
                        dataRight = [
                            ["Number of production: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfProduction'])],
                            ["One time production cost: USD " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['ProductionCostUsd'])],
                            ["Card Rate: " +
                                "Final Rates To Be Confirmed"],
                            ["Discounted rate: " +
                            "Final Rates To Be Confirmed"],
                            ["Total number of months required: " +
                                str(CampaignPlanDetail['MediaAdTypeDetail']['NoOfMonths'])],
                            ["Sub Total: " +
                            "Final Rates To Be Confirmed"]
                            # Add more rows as needed
                        ]

            column_padding = Inches(0.2)
            # Iterate through the data and populate table cells
            for row in range(min(len(dataLeft), rows)):
                for col in range(cols):
                    # Start from row 1 (excluding header row)
                    cell = table.cell(row, col)
                    cell.text = dataLeft[row][col]
                    cell.text_frame.paragraphs[0].font.size = Pt(16)
                    cell.margin_left = column_padding
                    cell.margin_right = column_padding
                    cell.margin_top = column_padding
                    cell.margin_bottom = column_padding
                    paragraph = text_frame.paragraphs[0]
                    paragraph.alignment = PP_ALIGN.CENTER
                    paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                    paragraph.font.size = Pt(12)

            row_colors = [RGBColor(240, 240, 240),
                        RGBColor(240, 240, 240)]
            for i, row in enumerate(table.rows):
                color_index = i % len(row_colors)
                for cell in row.cells:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = row_colors[color_index]

            # Define table dimensions and position
            rows = 6
            cols = 1
            left = slides_width/2 + Inches(0.5)
            top = Inches(3.3)
            width = int(slides_width/2 - Inches(1))
            height = Inches(2.5)
            # Add table to the slide
            table = slide.shapes.add_table(
                rows, cols, left, top, width, height).table
            table.style = "Table Grid"
            table.first_row = False

            column_padding = Inches(0.2)
            # Iterate through the data and populate table cells
            for row in range(min(len(dataRight), rows)):
                for col in range(cols):
                    # Start from row 1 (excluding header row)
                    cell = table.cell(row, col)
                    cell.text = dataRight[row][col]
                    cell.text_frame.paragraphs[0].font.size = Pt(16)
                    cell.margin_left = column_padding
                    cell.margin_right = column_padding
                    cell.margin_top = column_padding
                    cell.margin_bottom = column_padding
                    paragraph = text_frame.paragraphs[0]
                    paragraph.alignment = PP_ALIGN.CENTER
                    paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                    paragraph.font.size = Pt(12)

            row_colors = [RGBColor(240, 240, 240),
                        RGBColor(240, 240, 240)]
            for i, row in enumerate(table.rows):
                color_index = i % len(row_colors)
                for cell in row.cells:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = row_colors[color_index]

    # Slide 6: Total table slide
    itemdetail_slide_layout = presentation.slide_layouts[2]
    slide = presentation.slides.add_slide(itemdetail_slide_layout)

    # Remove main title
    title_shape = slide.shapes.title
    title_shape.text = " "

    # Check if the slide has a subtitle
    if len(slide.placeholders) > 1:
        # Remove the subtitle shape
        subtitle = slide.placeholders[1]
        slide.shapes._spTree.remove(subtitle._element)

    # item heading
    left = (slides_width/2)-Inches(3.8)/2
    top = Inches(0.5)
    width = Inches(3.8)
    height = Inches(1)
    table_head = slide.shapes.add_textbox(left, top, width, height)
    text_frame = table_head.text_frame
    itemHead = text_frame.add_paragraph()
    itemHead.text = "Estimated Media Investment"
    itemHead.font.size = Pt(22)

    # # Add table to the slide
    # table = slide.shapes.add_table(rows_per_slide + 1, cols, left, top, width, height).table

    # # Set table style
    # table.style = "Table Grid"

    # # Populate table cells
    # table.cell(0, 0).text = "Media Details"
    # table.cell(0, 1).text = "Total Number"
    # table.cell(0, 2).text = "No.of Production"
    # table.cell(0, 3).text = "Production Cost"
    # table.cell(0, 4).text = "Unit Cost"
    # table.cell(0, 5).text = "Discount Rate"
    # table.cell(0, 6).text = "Final Cost"
    EstimatedMediaInvestment = []
    for MediaInvestment in EstimateData['EstimatedMediaInvestment']:
        EstimatedMediaInvestment.append(
            [
                MediaInvestment['MediaDetails'],
                str(MediaInvestment['TotalNumber']),
                str(MediaInvestment['NoOfProduction']),
                str(MediaInvestment['ProductionCost']),
                str(MediaInvestment['UnitCost']),
                str(MediaInvestment['DiscountRate']),
                str(MediaInvestment['FinalCost']),
            ]
        )

    # Define table dimensions and position
    rows = len(EstimatedMediaInvestment)  # Maximum number of rows per slide
    cols = len(EstimatedMediaInvestment[0])
    left = Inches(0.5)
    top = Inches(1.7)
    width = int((slides_width) - Inches(1))
    height = Inches(2.5)
    rows_per_slide = 5
    column_padding = Inches(0.2)
    t_row = 0

    for row in range(rows):

        if row % rows_per_slide == 0:
            if row != 0:
                newslide_slide_layout = presentation.slide_layouts[2]
                new_slide = presentation.slides.add_slide(
                    newslide_slide_layout)
                top = Inches(0.5)
            else:
                new_slide = slide
                top = Inches(1.5)
            # Remove main title
            title_shape = new_slide.shapes.title
            title_shape.text = " "

            # Check if the slide has a subtitle
            if len(new_slide.placeholders) > 1:
                # Remove the subtitle shape
                subtitle = new_slide.placeholders[1]
                new_slide.shapes._spTree.remove(subtitle._element)

            remain_row = rows-row
            table_row_count = rows_per_slide+1
            if remain_row <= rows_per_slide:
                table_row_count = remain_row+1
            new_table = new_slide.shapes.add_table(
                table_row_count, cols, left, top, width, height).table

            # Set table style
            new_table.style = "Table Grid"

            new_table.cell(0, 0).text = "Media Details"
            new_table.cell(0, 1).text = "Total Number"
            new_table.cell(0, 2).text = "No.of Production"
            new_table.cell(0, 3).text = "Production Cost"
            new_table.cell(0, 4).text = "Unit Cost"
            new_table.cell(0, 5).text = "Discount Rate"
            new_table.cell(0, 6).text = "Final Cost"
            t_row = 0
        t_row = row % rows_per_slide
        for col in range(cols):
            cell = new_table.cell(t_row+1, col)
            if col==6:
                if EstimatedMediaInvestment[row][col]=="0.00":
                    cell.text ="Final rates to be confirmed"
                else:
                    cell.text = EstimatedMediaInvestment[row][col]
            else:
                if EstimatedMediaInvestment[row][col]=="0.00":
                    cell.text = "-"
                else:
                    cell.text = EstimatedMediaInvestment[row][col]
            

            cell.text_frame.paragraphs[0].font.size = Pt(16)
            cell.margin_left = column_padding
            cell.margin_right = column_padding
            cell.margin_top = column_padding
            cell.margin_bottom = column_padding
            paragraph = text_frame.paragraphs[0]
            paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        # print(t_row)

        column_widths = [Inches(3.8), Inches(1.2), Inches(
            1.5), Inches(2), Inches(1.7), Inches(1.9), Inches(1.9)]
        for i, column_width in enumerate(column_widths):
            column = new_table.columns[i]
            column.width = column_width
        # new_table.first_row = False
        header_row = new_table.rows[0]
        for cell in header_row.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(210, 210, 210)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        row_colors = [RGBColor(240, 240, 240),
                      RGBColor(240, 240, 240)]
        for i, row in enumerate(new_table.rows):
            color_index = i % len(row_colors)
            for cell in row.cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_colors[color_index]

    if t_row > 3:
        newslide_slide_layout = presentation.slide_layouts[2]
        new_slide = presentation.slides.add_slide(newslide_slide_layout)
        top = Inches(0.5)
    else:
        top = Inches(6.3)
    # Remove main title
    title_shape = new_slide.shapes.title
    title_shape.text = " "

    # Check if the slide has a subtitle
    if len(new_slide.placeholders) > 1:
        # Remove the subtitle shape
        subtitle = new_slide.placeholders[1]
        new_slide.shapes._spTree.remove(subtitle._element)

    # Define table dimensions and position
    rows = 4  # Maximum number of rows per slide
    cols = 1
    left = Inches(0.5)
    width = int((slides_width) - Inches(1))
    height = Inches(2.5)

    # Add table to the slide
    table = new_slide.shapes.add_table(
        rows, cols, left, top, width, height).table

    # Set table style
    table.style = "Table Grid"
    table.first_row = False
    # Populate table cells
    if(EstimateData['PlanDetail']['SubTotal']!="0.00"):
        table.cell(0, 0).text = "Sub Total in US Dollar: "+str(EstimateData['PlanDetail']['SubTotal'])
        table.cell(1, 0).text = "5% UAE VAT: "+str(EstimateData['PlanDetail']['VatAmount'])
        table.cell(2, 0).text = "Total in US Dollar: "+str(EstimateData['PlanDetail']['Total'])
        # table.cell(3, 0).text = EstimateData['PlanDetail']['TotalAmountInWords']
        cell_text = str(EstimateData['PlanDetail']['TotalAmountInWords'])
    else:
        table.cell(0, 0).text = "Sub Total in US Dollar: "+"Final rates to be confirmed"
        table.cell(1, 0).text = "5% UAE VAT: "+str(EstimateData['PlanDetail']['VatAmount'])
        table.cell(2, 0).text = "Total in US Dollar: "+"Final rates to be confirmed"
        # table.cell(3, 0).text = EstimateData['PlanDetail']['TotalAmountInWords']
        cell_text = "Final rates to be confirmed"
    capitalized_text = cell_text.title()
    table.cell(3, 0).text = capitalized_text
    
    # Iterate through the data and populate table cells
    column_padding = Inches(0.1)
    # Customize table cell formatting
    for row in range(rows):
        for col in range(cols):
            cell = table.cell(row, col)
            cell.text_frame.paragraphs[0].font.size = Pt(16)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            cell.margin_left = column_padding
            cell.margin_right = column_padding
            cell.margin_top = column_padding
            cell.margin_bottom = column_padding
            paragraph = text_frame.paragraphs[0]
            paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    header_row = table.rows[0]
    for cell in header_row.cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(220, 220, 220)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    row_colors = [RGBColor(240, 240, 240),
                  RGBColor(240, 240, 240)]
    for i, row in enumerate(table.rows):
        color_index = i % len(row_colors)
        for cell in row.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_colors[color_index]

    # Slide 7: Total table slide
    itemdetail_slide_layout = presentation.slide_layouts[2]
    slide = presentation.slides.add_slide(itemdetail_slide_layout)

    # Remove main title
    title_shape = slide.shapes.title
    title_shape.text = " "

    # Check if the slide has a subtitle
    if len(slide.placeholders) > 1:
        # Remove the subtitle shape
        subtitle = slide.placeholders[1]
        slide.shapes._spTree.remove(subtitle._element)

    # item heading
    left = (slides_width/2)-Inches(6.3)/2
    top = Inches(0.5)
    width = Inches(6.3)
    height = Inches(1)
    table_head = slide.shapes.add_textbox(left, top, width, height)
    text_frame = table_head.text_frame
    itemHead = text_frame.add_paragraph()
    itemHead.text = "Your campaign activation is just a few steps away. "
    itemHead.font.size = Pt(22)

    # Box 1 /////////////////
    background_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.75),
        Inches(2.5),
        Inches(3),
        Inches(4.5),
    )
    background_shape.fill.solid()
    background_shape.fill.fore_color.rgb = RGBColor(
        255, 255, 255)  # Set background color
    background_shape.shadow.inherit = False
    background_shape.shadow.visible = True

    # Add a heading
    heading_text_box = slide.shapes.add_textbox(
        left=Inches(0.75),
        top=Inches(3.8),
        width=Inches(3),
        height=Inches(1),
    )
    heading_text_frame = heading_text_box.text_frame
    heading_paragraph = heading_text_frame.add_paragraph()
    heading_paragraph.text = "Send For Approval"
    heading_paragraph.font.bold = True
    heading_paragraph.font.size = Pt(18)
    heading_paragraph.alignment = PP_ALIGN.CENTER

    # Add an image
    image_path = "images/estimate-icon01.png"
    image_left = Inches(1.75)
    image_top = Inches(3)
    image_width = Inches(1)
    image_height = Inches(1)
    slide.shapes.add_picture(image_path, image_left,
                             image_top, image_width, image_height)

    # Add a description
    description_text_box = slide.shapes.add_textbox(
        left=Inches(0.76),
        top=Inches(4.5),
        width=Inches(2.8),
        height=Inches(2),
    )
    description_text_frame = description_text_box.text_frame
    description_paragraph = description_text_frame.add_paragraph()
    description_paragraph.text = "Our support team will review your campaign plan and update you on the availability of selected dates. "
    description_paragraph.font.size = Pt(18)
    description_paragraph.alignment = PP_ALIGN.CENTER
    description_text_frame.word_wrap = True

    # Box 2 /////////////////
    background_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4.25),
        Inches(2.5),
        Inches(3),
        Inches(4.5),
    )
    background_shape.fill.solid()
    background_shape.fill.fore_color.rgb = RGBColor(
        255, 255, 255)  # Set background color
    background_shape.shadow.inherit = False
    background_shape.shadow.visible = True

    # Add a heading
    heading_text_box = slide.shapes.add_textbox(
        left=Inches(4.25),
        top=Inches(3.8),
        width=Inches(3),
        height=Inches(1),
    )
    heading_text_frame = heading_text_box.text_frame
    heading_paragraph = heading_text_frame.add_paragraph()
    heading_paragraph.text = "Review Media Plan"
    heading_paragraph.font.bold = True
    heading_paragraph.font.size = Pt(18)
    heading_paragraph.alignment = PP_ALIGN.CENTER

    # Add an image
    image_path = "images/estimate-icon02.png"
    image_left = Inches(5.25)
    image_top = Inches(3)
    image_width = Inches(1)
    image_height = Inches(1)
    slide.shapes.add_picture(image_path, image_left,
                             image_top, image_width, image_height)

    # Add a description
    description_text_box = slide.shapes.add_textbox(
        left=Inches(4.26),
        top=Inches(4.5),
        width=Inches(2.8),
        height=Inches(2),
    )
    description_text_frame = description_text_box.text_frame
    description_paragraph = description_text_frame.add_paragraph()
    description_paragraph.text = "You can review the plan and proceed to make payments based on the proforma invoice."
    description_paragraph.font.size = Pt(18)
    description_paragraph.alignment = PP_ALIGN.CENTER
    description_text_frame.word_wrap = True

    # Box 3 /////////////////
    background_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7.75),
        Inches(2.5),
        Inches(3),
        Inches(4.5),
    )
    background_shape.fill.solid()
    background_shape.fill.fore_color.rgb = RGBColor(
        255, 255, 255)  # Set background color
    background_shape.shadow.inherit = False
    background_shape.shadow.visible = True

    # Add a heading
    heading_text_box = slide.shapes.add_textbox(
        left=Inches(7.75),
        top=Inches(3.8),
        width=Inches(3),
        height=Inches(1),
    )
    heading_text_frame = heading_text_box.text_frame
    heading_paragraph = heading_text_frame.add_paragraph()
    heading_paragraph.text = "Booking"
    heading_paragraph.font.bold = True
    heading_paragraph.font.size = Pt(18)
    heading_paragraph.alignment = PP_ALIGN.CENTER

    # Add an image
    image_path = "images/estimate-icon03.png"
    image_left = Inches(8.75)
    image_top = Inches(3)
    image_width = Inches(1)
    image_height = Inches(1)
    slide.shapes.add_picture(image_path, image_left,
                             image_top, image_width, image_height)

    # Add a description
    description_text_box = slide.shapes.add_textbox(
        left=Inches(7.76),
        top=Inches(4.5),
        width=Inches(2.8),
        height=Inches(2),
    )
    description_text_frame = description_text_box.text_frame
    description_paragraph = description_text_frame.add_paragraph()
    description_paragraph.text = "Our team will get in touch to help you with an easy booking of your selected media options. "
    description_paragraph.font.size = Pt(18)
    description_paragraph.alignment = PP_ALIGN.CENTER
    description_text_frame.word_wrap = True

    # Box 4 /////////////////
    background_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(11.25),
        Inches(2.5),
        Inches(3),
        Inches(4.5),
    )
    background_shape.fill.solid()
    background_shape.fill.fore_color.rgb = RGBColor(
        255, 255, 255)  # Set background color
    background_shape.shadow.inherit = False
    background_shape.shadow.visible = True

    # Add a heading
    heading_text_box = slide.shapes.add_textbox(
        left=Inches(11.25),
        top=Inches(3.8),
        width=Inches(3),
        height=Inches(1),
    )
    heading_text_frame = heading_text_box.text_frame
    heading_paragraph = heading_text_frame.add_paragraph()
    heading_paragraph.text = "Launch"
    heading_paragraph.font.bold = True
    heading_paragraph.font.size = Pt(18)
    heading_paragraph.alignment = PP_ALIGN.CENTER

    # Add an image
    image_path = "images/estimate-icon04.png"
    image_left = Inches(12.25)
    image_top = Inches(3)
    image_width = Inches(1)
    image_height = Inches(1)
    slide.shapes.add_picture(image_path, image_left,
                             image_top, image_width, image_height)

    # Add a description
    description_text_box = slide.shapes.add_textbox(
        left=Inches(11.26),
        top=Inches(4.5),
        width=Inches(2.8),
        height=Inches(2),
    )
    description_text_frame = description_text_box.text_frame
    description_paragraph = description_text_frame.add_paragraph()
    description_paragraph.text = "We will ensure a successful launch of your campaign and provide you with post-campaign support."
    description_paragraph.font.size = Pt(18)
    description_paragraph.alignment = PP_ALIGN.CENTER
    description_text_frame.word_wrap = True

    # Slide 8: Total table slide
    itemdetail_slide_layout = presentation.slide_layouts[2]
    slide = presentation.slides.add_slide(itemdetail_slide_layout)

    # Remove main title
    title_shape = slide.shapes.title
    title_shape.text = " "

    # Check if the slide has a subtitle
    if len(slide.placeholders) > 1:
        # Remove the subtitle shape
        subtitle = slide.placeholders[1]
        slide.shapes._spTree.remove(subtitle._element)

    background_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        slides_width,
        Inches(3.5),
    )
    background_shape.fill.solid()
    background_shape.fill.fore_color.rgb = RGBColor(
        239, 239, 239)  # Set background color
    background_shape.shadow.inherit = False
    background_shape.shadow.visible = True

    # Add a heading
    heading_text_box = slide.shapes.add_textbox(
        left = Inches(0),
        top=Inches(1.0),
        width=slides_width,
        height=Inches(.5),
    )
    heading_text_frame = heading_text_box.text_frame
    heading_paragraph = heading_text_frame.add_paragraph()
    heading_paragraph.text = "Thank You"
    heading_paragraph.font.bold = True
    heading_paragraph.font.size = Pt(50)
    heading_paragraph.alignment = PP_ALIGN.CENTER

    left = Inches(0.5)
    top = Inches(3.5)
    width = Inches(14)
    height = Inches(1)
    table_head = slide.shapes.add_textbox(left, top, width, height)
    text_frame = table_head.text_frame
    itemHead = text_frame.add_paragraph()
    itemHead.text = " Masscom Global FZE, Office No: 220, One UAQ Building, UAQ Free Trade Zone, Umm Al Quwain, United Arab Emirates. "
    itemHead.font.size = Pt(18)
    itemHead.alignment = PP_ALIGN.CENTER
    itemHead.word_wrap = True

    # Add an image
    image_path = "images/mail.png"
    image_left = (slides_width/2)-Inches(5)/2 - Inches(0.4)
    image_top = Inches(4.4)
    image_width = Inches(0.4)
    image_height = Inches(0.4)
    slide.shapes.add_picture(image_path, image_left,
                             image_top, image_width, image_height)

    left = (slides_width/2)-Inches(5)/2
    top = Inches(4.1)
    width = Inches(2.5)
    height = Inches(1)
    table_head = slide.shapes.add_textbox(left, top, width, height)
    text_frame = table_head.text_frame
    itemHead = text_frame.add_paragraph()
    itemHead.text = "cs@globalmediakit.com"
    itemHead.font.size = Pt(18)
    itemHead.alignment = PP_ALIGN.CENTER
    itemHead.word_wrap = True

    # Add an image
    image_path = "images/phone.png"
    image_left = ((slides_width/2)-Inches(5)/2) + Inches(2.9)
    image_top = Inches(4.4)
    image_width = Inches(0.4)
    image_height = Inches(0.4)
    slide.shapes.add_picture(image_path, image_left,
                             image_top, image_width, image_height)

    left = ((slides_width/2)-Inches(5)/2) + Inches(2.9)
    top = Inches(4.1)
    width = Inches(2.5)
    height = Inches(1)
    table_head = slide.shapes.add_textbox(left, top, width, height)
    text_frame = table_head.text_frame
    itemHead = text_frame.add_paragraph()
    itemHead.text = "+971-43827760"
    itemHead.font.size = Pt(18)
    itemHead.alignment = PP_ALIGN.CENTER
    itemHead.word_wrap = True



    left = Inches(0)
    top = Inches(5)
    width = slides_width
    height = Inches(1)
    table_head = slide.shapes.add_textbox(left, top, width, height)
    text_frame = table_head.text_frame
    supportHead = text_frame.add_paragraph()
    # supportHead.text = "Need more support, schedule a meeting with us"
    supportHead.font.size = Pt(20)
    supportHead.alignment = PP_ALIGN.CENTER
    supportHead.word_wrap = True

    runsupportHead = supportHead.add_run()
    runsupportHead.text = "Need more support, schedule a meeting with us"
    # # Add a hyperlink to the text
    supportHeadhyperlink = runsupportHead.hyperlink
    supportHeadhyperlink.address = "https://calendly.com/gmksupport/global-media-kit-campaign-support"



    left = Inches(0.5)
    top = Inches(5.9)
    width = Inches(14)
    height = Inches(1)
    table_head = slide.shapes.add_textbox(left, top, width, height)
    text_frame = table_head.text_frame
    itemHead = text_frame.add_paragraph()
    itemHead.text = "Memberships"
    itemHead.font.size = Pt(20)
    itemHead.alignment = PP_ALIGN.CENTER
    itemHead.word_wrap = True

    background_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(6.6),
        Inches(14),
        Inches(2),
    )
    background_shape.fill.solid()
    background_shape.fill.fore_color.rgb = RGBColor(
        239, 239, 239)  # Set background color
    background_shape.shadow.inherit = False
    background_shape.shadow.visible = True

    # Add an image
    image_path = "images/h1.png"
    image_left = Inches(0.75)
    image_top = Inches(6.9)
    image_width = Inches(2.3)
    image_height = Inches(1.4)
    slide.shapes.add_picture(image_path, image_left,
                             image_top, image_width, image_height)

    # Add an image
    image_path = "images/h2.png"
    image_left = Inches(3.55)
    image_top = Inches(6.9)
    image_width = Inches(2.3)
    image_height = Inches(1.4)
    slide.shapes.add_picture(image_path, image_left,
                             image_top, image_width, image_height)

    # Add an image
    image_path = "images/h3.png"
    image_left = Inches(6.35)
    image_top = Inches(6.9)
    image_width = Inches(2.3)
    image_height = Inches(1.4)
    slide.shapes.add_picture(image_path, image_left,
                             image_top, image_width, image_height)

    # Add an image
    image_path = "images/h4.png"
    image_left = Inches(9.15)
    image_top = Inches(6.9)
    image_width = Inches(2.3)
    image_height = Inches(1.4)
    slide.shapes.add_picture(image_path, image_left,
                             image_top, image_width, image_height)

    # Add an image
    image_path = "images/h5.png"
    image_left = Inches(11.95)
    image_top = Inches(6.9)
    image_width = Inches(2.3)
    image_height = Inches(1.4)
    slide.shapes.add_picture(image_path, image_left,
                             image_top, image_width, image_height)

    # Save the presentation 
    isExist = os.path.exists('media/estimate/excel/')
    if not isExist:
        os.makedirs('media/estimate/excel/')
    presentation.save('media/estimate/excel/'+EstimateData['PlanDetail']['PlanName']+".pptx")

    return JsonResponse('/media/estimate/excel/'+EstimateData['PlanDetail']['PlanName']+".pptx", safe=False)
